from pptx import Presentation
from pptx.util import Inches

# Create a presentation object
prs = Presentation()

# Define slide titles and content
slides_content = [
    ("MediSortix: AI-Powered ER Patient Prioritization System",
     "Improving Patient Outcomes with AI-Driven Prioritization\nPresented by: [Your Name]\nDate: [Date]"),

    ("Introduction",
     "Introduce yourself and your role.\nOverview of the presentation: Introducing MediSortix, its functionalities, and benefits."),

    ("The Problem",
     "Challenges in ERs with patient prioritization.\nHigh stakes in correctly prioritizing patients to reduce morbidity and mortality."),

    ("The Solution - MediSortix",
     "Introducing MediSortix as a comprehensive solution.\nCombining AI, chatbots, and SHAP explainers to prioritize patients effectively."),

    ("System Workflow",
     "Visual diagram of the system workflow.\n1. Patient arrives at the ER.\n2. Chatbot interacts with the patient to gather information.\n3. Data processed by the AI model.\n4. SHAP explainer calculates feature importance.\n5. Patients reordered based on SHAP values."),

    ("Chatbot Interaction",
     "Chatbot's role in collecting patient information.\nKey features: conversational interface, comprehensive medical history, and baseline parameters."),

    ("AI Model and SHAP Explainer",
     "Explain the AI model (LightGBM) and its purpose in predicting patient outcomes.\nDescribe SHAP explainer and feature importance calculation."),

    ("Patient Reordering",
     "Show how the system reorders patients based on SHAP values.\nDescending order of risk, prioritizing those with the highest risk of death."),

    ("Benefits of MediSortix",
     "Benefits: improved patient outcomes, efficient prioritization, reduced ER congestion, and data-driven decision-making.\nPotential impact on reducing morbidity and mortality rates."),

    ("Conclusion and Future Work",
     "Summarize key points.\nDiscuss potential future enhancements: integration with other hospital systems, continuous learning, expanding to other departments."),

    ("Q&A",
     "Open the floor for questions.\nEncourage interaction and address queries regarding the system."),
]

# Add slides to the presentation
for title, content in slides_content:
    slide_layout = prs.slide_layouts[1]  # Use the 'Title and Content' layout
    slide = prs.slides.add_slide(slide_layout)
    title_placeholder = slide.shapes.title
    content_placeholder = slide.placeholders[1]

    title_placeholder.text = title
    content_placeholder.text = content

# Save the presentation
prs.save("MediSortix_Presentation.pptx")